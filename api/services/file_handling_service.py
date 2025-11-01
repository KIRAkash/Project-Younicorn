"""File handling service for extracting text from various file formats."""

import logging
import tempfile
import os
import shutil
from typing import Dict, List, Optional
from google.cloud import speech_v1
from google.cloud import storage
import subprocess
import mimetypes

# Document processing libraries
try:
    import PyPDF2
    from docx import Document
    from pptx import Presentation
    import pandas as pd
except ImportError:
    PyPDF2 = None
    Document = None
    Presentation = None
    pd = None

logger = logging.getLogger(__name__)

# Import cache service (will be initialized after class definition)
file_content_cache_service = None


class FileHandlingService:
    """Service for extracting text content from various file formats."""
    
    def __init__(self):
        """Initialize the file handling service."""
        self.speech_client = speech_v1.SpeechClient()
        self.storage_client = storage.Client()
        
    def _download_from_gcs(self, gcs_uri: str, local_path: str) -> bool:
        """Download a file from GCS to local path."""
        try:
            # Parse GCS URI (gs://bucket/path)
            if not gcs_uri.startswith("gs://"):
                logger.error(f"Invalid GCS URI: {gcs_uri}")
                return False
                
            parts = gcs_uri[5:].split("/", 1)
            if len(parts) != 2:
                logger.error(f"Invalid GCS URI format: {gcs_uri}")
                return False
                
            bucket_name, blob_name = parts
            bucket = self.storage_client.bucket(bucket_name)
            blob = bucket.blob(blob_name)
            
            blob.download_to_filename(local_path)
            logger.info(f"Downloaded {gcs_uri} to {local_path}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to download from GCS {gcs_uri}: {e}")
            return False
    
    def _extract_audio_from_video_gcs(self, video_gcs_uri: str, audio_gcs_uri: str) -> bool:
        """Extract audio from video file using ffmpeg (download, process, upload)."""
        temp_dir = None
        try:
            # Create temporary directory
            temp_dir = tempfile.mkdtemp()
            video_path = os.path.join(temp_dir, "video.mp4")
            audio_path = os.path.join(temp_dir, "audio.wav")
            
            # Download video from GCS
            if not self._download_from_gcs(video_gcs_uri, video_path):
                return False
            
            # Use ffmpeg to extract audio
            command = [
                'ffmpeg',
                '-i', video_path,
                '-vn',  # No video
                '-acodec', 'pcm_s16le',  # Linear PCM 16-bit
                '-ar', '16000',  # 16kHz sample rate
                '-ac', '1',  # Mono
                audio_path,
                '-y'  # Overwrite output file
            ]
            
            result = subprocess.run(
                command,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=300  # 5 minute timeout
            )
            
            if result.returncode != 0:
                logger.error(f"ffmpeg failed: {result.stderr.decode()}")
                return False
            
            logger.info(f"Successfully extracted audio from video")
            
            # Upload audio to GCS
            parts = audio_gcs_uri[5:].split("/", 1)
            if len(parts) != 2:
                logger.error(f"Invalid GCS URI format: {audio_gcs_uri}")
                return False
            
            bucket_name, blob_name = parts
            bucket = self.storage_client.bucket(bucket_name)
            blob = bucket.blob(blob_name)
            blob.upload_from_filename(audio_path)
            logger.info(f"Uploaded audio to {audio_gcs_uri}")
            
            return True
                
        except subprocess.TimeoutExpired:
            logger.error(f"Audio extraction timed out")
            return False
        except FileNotFoundError:
            logger.error("ffmpeg not found. Please install ffmpeg.")
            return False
        except Exception as e:
            logger.error(f"Failed to extract audio: {e}")
            return False
        finally:
            # Clean up temporary directory
            if temp_dir and os.path.exists(temp_dir):
                try:
                    shutil.rmtree(temp_dir)
                    logger.debug(f"Cleaned up temporary directory: {temp_dir}")
                except Exception as e:
                    logger.warning(f"Failed to clean up temp directory {temp_dir}: {e}")
    
    def _transcribe_audio_with_speech_api(self, gcs_uri: str, content_type: str) -> Optional[str]:
        """Transcribe audio file using Google Speech-to-Text API."""
        try:
            # Configure audio settings
            audio = speech_v1.RecognitionAudio(uri=gcs_uri)
            
            # Determine encoding and sample rate from content type
            encoding = speech_v1.RecognitionConfig.AudioEncoding.LINEAR16
            sample_rate = 16000
            
            if 'flac' in content_type.lower():
                encoding = speech_v1.RecognitionConfig.AudioEncoding.FLAC
                sample_rate = None  # FLAC has embedded sample rate
            elif 'mp3' in content_type.lower() or 'mpeg' in content_type.lower():
                encoding = speech_v1.RecognitionConfig.AudioEncoding.MP3
                sample_rate = None  # MP3 has variable sample rate, let API auto-detect
            elif 'webm' in content_type.lower():
                encoding = speech_v1.RecognitionConfig.AudioEncoding.WEBM_OPUS
                sample_rate = None  # WEBM Opus has embedded sample rate (usually 48000)
            elif 'ogg' in content_type.lower() or 'opus' in content_type.lower():
                encoding = speech_v1.RecognitionConfig.AudioEncoding.OGG_OPUS
                sample_rate = None  # OGG Opus has embedded sample rate
            elif 'wav' in content_type.lower():
                encoding = speech_v1.RecognitionConfig.AudioEncoding.LINEAR16
                sample_rate = 16000
            
            # Build config with optional sample_rate_hertz
            # Note: Using "video" model which works better for general audio/video content
            # "video" model only supports en-US language code
            config_params = {
                "encoding": encoding,
                "language_code": "en-US",  # Primary language - English (US)
                "enable_automatic_punctuation": True,
                "enable_word_time_offsets": True,
                "model": "video",  # Better for general audio/video transcription
            }
            
            # Only add sample_rate_hertz if we have a specific rate
            if sample_rate:
                config_params["sample_rate_hertz"] = sample_rate
            
            config = speech_v1.RecognitionConfig(**config_params)
            
            # Use long-running recognize for all GCS files (required for files > 1 minute)
            # Note: For GCS URIs, we must use long_running_recognize
            logger.info(f"Starting transcription for {gcs_uri} with encoding={encoding}")
            logger.info(f"Config: language={config.language_code}, sample_rate={'auto-detect' if not sample_rate else sample_rate}Hz")
            
            operation = self.speech_client.long_running_recognize(
                config=config,
                audio=audio
            )
            
            logger.info("Waiting for transcription to complete...")
            response = operation.result(timeout=600)  # 10 minute timeout
            
            # Log response details for debugging
            logger.info(f"Transcription response received. Results count: {len(response.results)}")
            
            # Combine all transcripts (only non-empty ones)
            transcript_parts = []
            for idx, result in enumerate(response.results):
                logger.debug(f"Result {idx}: {len(result.alternatives)} alternatives")
                if result.alternatives:
                    transcript = result.alternatives[0].transcript.strip()
                    confidence = result.alternatives[0].confidence if hasattr(result.alternatives[0], 'confidence') else 0
                    if transcript:  # Only add non-empty transcripts
                        logger.debug(f"  Transcript: {transcript[:100]}... (confidence: {confidence})")
                        transcript_parts.append(transcript)
                    else:
                        logger.debug(f"  Empty transcript (confidence: {confidence})")
            
            if not transcript_parts:
                logger.warning(f"No transcription results found with video model. Trying with minimal config...")
                
                # Fallback: Try with minimal configuration
                try:
                    minimal_config = speech_v1.RecognitionConfig(
                        encoding=encoding,
                        language_code="en-US",
                        enable_automatic_punctuation=True,
                    )
                    if sample_rate:
                        minimal_config.sample_rate_hertz = sample_rate
                    
                    operation2 = self.speech_client.long_running_recognize(
                        config=minimal_config,
                        audio=audio
                    )
                    response2 = operation2.result(timeout=600)
                    
                    logger.info(f"Fallback transcription response: {len(response2.results)} results")
                    
                    for result in response2.results:
                        if result.alternatives:
                            transcript_parts.append(result.alternatives[0].transcript)
                    
                    if not transcript_parts:
                        logger.warning(f"No transcription results found even with fallback for {gcs_uri}. The audio may be silent, too short, or in an unsupported format.")
                        return None
                        
                except Exception as fallback_error:
                    logger.error(f"Fallback transcription also failed: {fallback_error}")
                    return None
            
            full_transcript = " ".join(transcript_parts)
            logger.info(f"Transcription completed. Length: {len(full_transcript)} characters, {len(transcript_parts)} segments")
            
            return full_transcript
            
        except Exception as e:
            logger.error(f"Failed to transcribe audio {gcs_uri}: {e}")
            return None
    
    def _extract_text_from_pdf(self, file_path: str) -> Optional[str]:
        """Extract text from PDF file."""
        if PyPDF2 is None:
            logger.error("PyPDF2 not installed. Cannot extract PDF text.")
            return None
            
        try:
            text_parts = []
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text:
                        text_parts.append(f"[Page {page_num + 1}]\n{text}")
            
            full_text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(full_text)} characters from PDF")
            return full_text
            
        except Exception as e:
            logger.error(f"Failed to extract text from PDF: {e}")
            return None
    
    def _extract_text_from_docx(self, file_path: str) -> Optional[str]:
        """Extract text from DOCX file."""
        if Document is None:
            logger.error("python-docx not installed. Cannot extract DOCX text.")
            return None
            
        try:
            doc = Document(file_path)
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            full_text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(full_text)} characters from DOCX")
            return full_text
            
        except Exception as e:
            logger.error(f"Failed to extract text from DOCX: {e}")
            return None
    
    def _extract_text_from_pptx(self, file_path: str) -> Optional[str]:
        """Extract text from PPTX file."""
        if Presentation is None:
            logger.error("python-pptx not installed. Cannot extract PPTX text.")
            return None
            
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"[Slide {slide_num + 1}]"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += f"\n{shape.text}"
                
                if slide_text != f"[Slide {slide_num + 1}]":
                    text_parts.append(slide_text)
            
            full_text = "\n\n".join(text_parts)
            logger.info(f"Extracted {len(full_text)} characters from PPTX")
            return full_text
            
        except Exception as e:
            logger.error(f"Failed to extract text from PPTX: {e}")
            return None
    
    def _extract_text_from_txt(self, file_path: str) -> Optional[str]:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text = file.read()
            
            logger.info(f"Extracted {len(text)} characters from text file")
            return text
            
        except Exception as e:
            logger.error(f"Failed to extract text from text file: {e}")
            return None
    
    def _extract_text_from_csv(self, file_path: str) -> Optional[str]:
        """Extract text from CSV file."""
        if pd is None:
            logger.error("pandas not installed. Cannot extract CSV text.")
            return None
            
        try:
            df = pd.read_csv(file_path)
            
            # Convert DataFrame to readable text format
            text_parts = []
            text_parts.append(f"CSV File with {len(df)} rows and {len(df.columns)} columns\n")
            text_parts.append(f"Columns: {', '.join(df.columns)}\n")
            text_parts.append("\nData Preview (first 10 rows):\n")
            text_parts.append(df.head(10).to_string())
            
            # Add summary statistics for numeric columns
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                text_parts.append("\n\nSummary Statistics:\n")
                text_parts.append(df[numeric_cols].describe().to_string())
            
            full_text = "\n".join(text_parts)
            logger.info(f"Extracted {len(full_text)} characters from CSV")
            return full_text
            
        except Exception as e:
            logger.error(f"Failed to extract text from CSV: {e}")
            return None
    
    def _extract_text_from_excel(self, file_path: str) -> Optional[str]:
        """Extract text from Excel file (xlsx, xls)."""
        if pd is None:
            logger.error("pandas not installed. Cannot extract Excel text.")
            return None
            
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text_parts = []
            
            text_parts.append(f"Excel File with {len(excel_file.sheet_names)} sheet(s): {', '.join(excel_file.sheet_names)}\n")
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                text_parts.append(f"\n{'='*60}")
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append(f"{'='*60}")
                text_parts.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
                text_parts.append(f"Columns: {', '.join(df.columns)}\n")
                text_parts.append("Data Preview (first 10 rows):\n")
                text_parts.append(df.head(10).to_string())
                
                # Add summary statistics for numeric columns
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    text_parts.append("\nSummary Statistics:\n")
                    text_parts.append(df[numeric_cols].describe().to_string())
            
            full_text = "\n".join(text_parts)
            logger.info(f"Extracted {len(full_text)} characters from Excel")
            return full_text
            
        except Exception as e:
            logger.error(f"Failed to extract text from Excel: {e}")
            return None
    
    async def extract_text_from_file(self, gcs_uri: str, content_type: str, filename: str) -> Optional[str]:
        """
        Extract text from a file based on its type.
        Checks cache first to avoid redundant processing.
        
        Args:
            gcs_uri: GCS URI of the file (gs://bucket/path)
            content_type: MIME type of the file
            filename: Original filename
            
        Returns:
            Extracted text content or None if extraction failed
        """
        logger.info(f"Extracting text from {filename} ({content_type}) at {gcs_uri}")
        
        # Check cache first
        global file_content_cache_service
        if file_content_cache_service is None:
            try:
                from .file_content_cache_service import file_content_cache_service as cache_svc
                file_content_cache_service = cache_svc
            except Exception as e:
                logger.warning(f"Could not import cache service: {e}")
        
        if file_content_cache_service:
            cached_content = file_content_cache_service.get_cached_content(gcs_uri, filename, content_type)
            if cached_content:
                logger.info(f"✓ Using CACHED content for {filename} (accessed {cached_content.get('cache_access_count', 0)} times)")
                if cached_content.get('processing_status') == 'success':
                    return cached_content.get('extracted_text')
                elif cached_content.get('processing_status') == 'failed':
                    logger.warning(f"Previously failed to extract from {filename}: {cached_content.get('error_message')}")
                    return None
                elif cached_content.get('processing_status') == 'empty':
                    logger.warning(f"Previously extracted empty content from {filename}")
                    return None
        
        logger.info(f"⚡ Processing NEW file: {filename}")
        
        # Process the file and cache the result
        extracted_text = None
        processing_status = "success"
        error_message = None
        
        try:
            # Handle video files - extract audio first, then transcribe
            if content_type.startswith('video/'):
                logger.info(f"Processing video file: {filename}")
                
                # Generate audio GCS URI
                audio_gcs_uri = gcs_uri.replace(os.path.splitext(filename)[1], '_audio.wav')
                
                # Extract audio from video and upload to GCS
                if not self._extract_audio_from_video_gcs(gcs_uri, audio_gcs_uri):
                    logger.warning("Failed to extract audio from video, attempting direct transcription")
                    extracted_text = self._transcribe_audio_with_speech_api(gcs_uri, content_type)
                else:
                    # Transcribe audio using GCS URI (no download needed)
                    extracted_text = self._transcribe_audio_with_speech_api(audio_gcs_uri, 'audio/wav')
            
            # Handle audio files - transcribe directly using GCS URI (no download needed)
            elif content_type.startswith('audio/'):
                logger.info(f"Processing audio file: {filename}")
                extracted_text = self._transcribe_audio_with_speech_api(gcs_uri, content_type)
            
            # Handle document files - need to download temporarily for text extraction
            else:
                temp_dir = None
                try:
                    # Create temporary directory
                    temp_dir = tempfile.mkdtemp()
                    local_path = os.path.join(temp_dir, filename)
                    
                    # Download file
                    if not self._download_from_gcs(gcs_uri, local_path):
                        processing_status = "failed"
                        error_message = "Failed to download file from GCS"
                        return None
                    
                    # Extract text based on file type
                    if content_type == 'application/pdf' or filename.lower().endswith('.pdf'):
                        extracted_text = self._extract_text_from_pdf(local_path)
                    
                    elif content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or filename.lower().endswith('.docx'):
                        extracted_text = self._extract_text_from_docx(local_path)
                    
                    elif content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or filename.lower().endswith('.pptx'):
                        extracted_text = self._extract_text_from_pptx(local_path)
                    
                    elif content_type == 'text/csv' or filename.lower().endswith('.csv'):
                        extracted_text = self._extract_text_from_csv(local_path)
                    
                    elif content_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel'] or filename.lower().endswith(('.xlsx', '.xls')):
                        extracted_text = self._extract_text_from_excel(local_path)
                    
                    elif content_type.startswith('text/') or filename.lower().endswith(('.txt', '.md')):
                        extracted_text = self._extract_text_from_txt(local_path)
                    
                    else:
                        logger.warning(f"Unsupported file type: {content_type} for {filename}")
                        processing_status = "failed"
                        error_message = f"Unsupported file type: {content_type}"
                        extracted_text = None
                    
                finally:
                    # Clean up temporary directory
                    if temp_dir and os.path.exists(temp_dir):
                        try:
                            shutil.rmtree(temp_dir)
                            logger.debug(f"Cleaned up temporary directory: {temp_dir}")
                        except Exception as e:
                            logger.warning(f"Failed to clean up temp directory {temp_dir}: {e}")
            
            # Determine processing status
            if extracted_text is None:
                if processing_status == "success":
                    processing_status = "failed"
                    error_message = error_message or "Text extraction returned None"
            elif len(extracted_text.strip()) == 0:
                processing_status = "empty"
                error_message = "No text content found in file"
            
        except Exception as e:
            logger.error(f"Error processing file {filename}: {e}")
            processing_status = "failed"
            error_message = str(e)
            extracted_text = None
        
        # Cache the result (success, failed, or empty)
        if file_content_cache_service:
            try:
                file_content_cache_service.cache_file_content(
                    gcs_uri=gcs_uri,
                    filename=filename,
                    content_type=content_type,
                    extracted_text=extracted_text,
                    processing_status=processing_status,
                    error_message=error_message
                )
            except Exception as cache_error:
                logger.warning(f"Failed to cache file content: {cache_error}")
        
        return extracted_text
    
    async def process_files(self, gcs_files: List[Dict]) -> List[Dict[str, str]]:
        """
        Process multiple files and extract text from each.
        Checks cache first to avoid redundant processing.
        
        Args:
            gcs_files: List of file info dicts with gcs_path, content_type, filename
            
        Returns:
            List of dicts with filename, extracted_text, and cached flag
        """
        attachments = []
        
        # Initialize cache service if not already done
        global file_content_cache_service
        if file_content_cache_service is None:
            try:
                from .file_content_cache_service import file_content_cache_service as cache_svc
                file_content_cache_service = cache_svc
            except Exception as e:
                logger.warning(f"Could not import cache service: {e}")
        
        for file_info in gcs_files:
            gcs_uri = file_info.get('gcs_path')
            content_type = file_info.get('content_type', '')
            filename = file_info.get('filename', 'unknown')
            
            if not gcs_uri:
                logger.warning(f"No GCS path for file: {filename}")
                continue
            
            # Check cache first
            cached = False
            cached_content = None
            if file_content_cache_service:
                cached_content = file_content_cache_service.get_cached_content(gcs_uri, filename, content_type)
                if cached_content and cached_content.get('processing_status') == 'success':
                    cached = True
            
            # Extract text (will use cache internally if available)
            extracted_text = await self.extract_text_from_file(gcs_uri, content_type, filename)
            
            if extracted_text and len(extracted_text.strip()) > 0:
                attachments.append({
                    "filename": filename,
                    "content_type": content_type,
                    "extracted_text": extracted_text,
                    "text_length": len(extracted_text),
                    "cached": cached  # Flag to indicate if content was from cache
                })
                cache_info = " (from cache)" if cached else ""
                logger.info(f"Successfully extracted {len(extracted_text)} characters from {filename}{cache_info}")
            else:
                error_msg = f"[No text could be extracted from {filename}. "
                if content_type.startswith('audio/') or content_type.startswith('video/'):
                    error_msg += "The audio may be silent, too short, or speech may not be clear enough for transcription.]"
                else:
                    error_msg += "The file may be empty, corrupted, or in an unsupported format.]"
                
                logger.warning(f"No text extracted from {filename}")
                attachments.append({
                    "filename": filename,
                    "content_type": content_type,
                    "extracted_text": error_msg,
                    "text_length": 0,
                    "cached": False  # Failed extractions are not cached
                })
        
        return attachments


# Global service instance
file_handling_service = FileHandlingService()
